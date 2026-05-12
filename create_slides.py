"""Generate agents4geos project slides."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree

# ── Colour palette ──────────────────────────────────────────────────────────
DARK_BG    = RGBColor(0x1A, 0x1A, 0x2E)   # deep navy
ACCENT1    = RGBColor(0x16, 0x21, 0x3E)   # mid navy (panel bg)
TEAL       = RGBColor(0x0F, 0xBC, 0xD4)   # bright teal (headings)
ORANGE     = RGBColor(0xF5, 0xA6, 0x23)   # amber (highlights)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
GREEN      = RGBColor(0x4C, 0xAF, 0x50)
PURPLE     = RGBColor(0xAB, 0x47, 0xBC)
RED_SOFT   = RGBColor(0xEF, 0x53, 0x50)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def set_bg(slide, color: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, l, t, w, h, fill_color, alpha=None):
    shape = slide.shapes.add_shape(
        pptx.enum.shapes.MSO_SHAPE_TYPE.AUTO_SHAPE if False else 1,
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_text(slide, text, l, t, w, h, font_size=18, bold=False,
             color=WHITE, align=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_para(tf, text, font_size=14, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False, space_before=0):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return p


prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]   # completely blank layout


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title / Overview
# ════════════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(blank)
set_bg(s1, DARK_BG)

# Top accent bar
add_rect(s1, 0, 0, 13.33, 0.12, TEAL)

# Bottom accent bar
add_rect(s1, 0, 7.38, 13.33, 0.12, TEAL)

# Left vertical accent stripe
add_rect(s1, 0, 0.12, 0.08, 7.26, ORANGE)

# Main title
add_text(s1, "Agents4GEOS", 0.3, 0.5, 8, 1.2,
         font_size=52, bold=True, color=TEAL, align=PP_ALIGN.LEFT)

# Subtitle
add_text(s1, "A multi-agent MCP server for natural-language GEOS simulation design",
         0.3, 1.6, 9, 0.7, font_size=22, color=LIGHT_GRAY)

# Divider line (thin rect)
add_rect(s1, 0.3, 2.35, 12.7, 0.03, ORANGE)

# Stats row — three panels
panels = [
    ("52", "MCP Tools",    TEAL,   2.0),
    ("11", "AI Agents",    ORANGE, 5.2),
    ("200+", "GEOS Examples\n(knowledge base)", GREEN, 8.4),
]
for num, label, col, left in panels:
    add_rect(s1, left, 2.55, 2.9, 1.7, ACCENT1)
    add_text(s1, num,   left+0.1, 2.65, 2.7, 0.85,
             font_size=44, bold=True, color=col, align=PP_ALIGN.CENTER)
    add_text(s1, label, left+0.1, 3.4,  2.7, 0.7,
             font_size=16, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# What it does
add_text(s1, "What it does", 0.3, 4.45, 5, 0.4,
         font_size=18, bold=True, color=ORANGE)
desc_box = s1.shapes.add_textbox(Inches(0.3), Inches(4.85), Inches(7.8), Inches(2.2))
dtf = desc_box.text_frame
dtf.word_wrap = True
items = [
    "• Converts plain-English simulation descriptions into validated GEOS XML",
    "• Handles single-phase / multiphase / CO₂-brine / thermal coupling",
    "• Covers the full workflow: mesh → fluids → XML assembly → validation → post-processing",
    "• Built on FastMCP + Claude Code slash commands",
]
for i, item in enumerate(items):
    if i == 0:
        p = dtf.paragraphs[0]
    else:
        p = dtf.add_paragraph()
    p.space_before = Pt(4)
    r = p.add_run()
    r.text = item
    r.font.size = Pt(16)
    r.font.color.rgb = WHITE

# Tech stack badge (right side)
add_rect(s1, 8.5, 4.45, 4.6, 2.7, ACCENT1)
add_text(s1, "Tech Stack", 8.65, 4.55, 4.3, 0.4,
         font_size=16, bold=True, color=TEAL)
stack_box = s1.shapes.add_textbox(Inches(8.65), Inches(4.95), Inches(4.3), Inches(2.0))
stf = stack_box.text_frame
stf.word_wrap = True
stack = [
    ("FastMCP",        ORANGE),
    ("geos-tui  (schema / XML)", LIGHT_GRAY),
    ("pyResToolbox  (fluid PVT)", LIGHT_GRAY),
    ("PyVista  (mesh / VTK)", LIGHT_GRAY),
    ("Claude Haiku / Sonnet / Opus", TEAL),
]
for i, (item, col) in enumerate(stack):
    p = stf.paragraphs[0] if i == 0 else stf.add_paragraph()
    p.space_before = Pt(3)
    r = p.add_run()
    r.text = f"  {item}"
    r.font.size = Pt(15)
    r.font.color.rgb = col

# Slide number
add_text(s1, "1 / 4", 12.8, 7.1, 0.4, 0.3, font_size=11,
         color=LIGHT_GRAY, align=PP_ALIGN.RIGHT)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Tool Architecture (52 tools in 6 modules)
# ════════════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(blank)
set_bg(s2, DARK_BG)
add_rect(s2, 0, 0, 13.33, 0.12, TEAL)
add_rect(s2, 0, 7.38, 13.33, 0.12, TEAL)
add_rect(s2, 0, 0.12, 0.08, 7.26, ORANGE)

add_text(s2, "Tool Architecture  —  52 MCP Tools in 6 Modules",
         0.25, 0.18, 13, 0.55, font_size=28, bold=True, color=TEAL)

modules = [
    # (module, count, colour, sample tools)
    ("Schema & Introspection",  7, TEAL,
     "list_elements · get_element_info\nget_children · get_required_attrs\nvalidate_path · find_cross_refs"),
    ("Fluid & Constitutive",   10, ORANGE,
     "compute_pvt · recommend_fluid_model\nget_viscosity · get_density\nbuild_co2_brine · build_dead_oil"),
    ("Mesh",                    8, GREEN,
     "create_box_mesh · create_vtk_mesh\nvisualise_mesh · mesh_statistics\nconvert_mesh · inspect_gmsh"),
    ("XML Assembly & Validation", 14, PURPLE,
     "create_document · add_element\nset_attribute · validate_xml\ncross_ref_check · save_document"),
    ("Post-Processing",         9, RED_SOFT,
     "list_vtk_fields · plot_field\ntime_series · compute_gradient\nextract_slice · summarise_output"),
    ("Pre-Processing / Utilities", 4, RGBColor(0xFF,0xB7,0x4D),
     "convert_units · expand_params\nresolve_includes · health_check"),
]

cols = 3
rows = 2
cell_w = 4.2
cell_h = 2.8
pad_x  = 0.12
pad_y  = 0.2
start_x = 0.2
start_y = 0.85

for idx, (name, count, col, tools) in enumerate(modules):
    c = idx % cols
    r = idx // cols
    lft = start_x + c * (cell_w + pad_x)
    top = start_y + r * (cell_h + pad_y)
    add_rect(s2, lft, top, cell_w, cell_h, ACCENT1)
    # Colour top strip
    add_rect(s2, lft, top, cell_w, 0.07, col)
    # Count badge
    add_rect(s2, lft + cell_w - 0.7, top + 0.1, 0.6, 0.5, col)
    add_text(s2, str(count), lft + cell_w - 0.7, top + 0.1, 0.6, 0.5,
             font_size=18, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)
    # Module name
    add_text(s2, name, lft + 0.15, top + 0.15, cell_w - 0.9, 0.55,
             font_size=14, bold=True, color=col)
    # Sample tools
    tb = s2.shapes.add_textbox(Inches(lft + 0.15), Inches(top + 0.7),
                                Inches(cell_w - 0.3), Inches(cell_h - 0.8))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(tools.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r2 = p.add_run()
        r2.text = line
        r2.font.size = Pt(12)
        r2.font.color.rgb = LIGHT_GRAY

add_text(s2, "2 / 4", 12.8, 7.1, 0.4, 0.3, font_size=11,
         color=LIGHT_GRAY, align=PP_ALIGN.RIGHT)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Agent Taxonomy & Orchestration
# ════════════════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(blank)
set_bg(s3, DARK_BG)
add_rect(s3, 0, 0, 13.33, 0.12, TEAL)
add_rect(s3, 0, 7.38, 13.33, 0.12, TEAL)
add_rect(s3, 0, 0.12, 0.08, 7.26, ORANGE)

add_text(s3, "Agent Taxonomy  &  Orchestration Flow",
         0.25, 0.18, 13, 0.55, font_size=28, bold=True, color=TEAL)

# Left: agent table
add_rect(s3, 0.2, 0.85, 7.1, 6.3, ACCENT1)
add_text(s3, "11 Specialised Agents", 0.35, 0.9, 6.8, 0.45,
         font_size=17, bold=True, color=ORANGE)

agents = [
    # agent, tier, model, role
    ("/geos",              "3", "Opus",   "Main orchestrator  (all tools)"),
    ("/geos:edit",         "2", "Sonnet", "XML modification & editing"),
    ("/geos:validate",     "2", "Sonnet", "Schema + physics validation"),
    ("/geos:fluids",       "2", "Sonnet", "Fluid PVT & constitutive models"),
    ("/geos:mesh",         "2", "Sonnet", "Mesh creation & visualisation"),
    ("/geos:relperm",      "2", "Sonnet", "Rel-perm / cap-pressure curves"),
    ("/geos:postprocess",  "2", "Sonnet", "VTK output analysis"),
    ("/geos:schema",       "1", "Haiku",  "Schema introspection (read-only)"),
    ("/geos:inspect",      "1", "Haiku",  "XML inspection (read-only)"),
    ("/geos:run",          "3", "Opus",   "Run simulation & diagnose errors"),
    ("/geos:curate-errors","1", "Haiku",  "Log runtime errors for learning"),
]

tier_cols = {"1": RGBColor(0x42,0xA5,0xF5), "2": ORANGE, "3": GREEN}
model_cols = {"Haiku": RGBColor(0x42,0xA5,0xF5), "Sonnet": ORANGE, "Opus": GREEN}

# Header row
hdr_y = 1.38
add_rect(s3, 0.25, hdr_y, 7.0, 0.32, RGBColor(0x0D,0x47,0xA1))
for txt, x, w in [("Agent", 0.3, 2.5), ("Tier", 2.85, 0.5),
                   ("Model", 3.4, 1.2), ("Role", 4.65, 2.5)]:
    add_text(s3, txt, x, hdr_y+0.02, w, 0.28,
             font_size=11, bold=True, color=WHITE)

for i, (agent, tier, model, role) in enumerate(agents):
    row_y = hdr_y + 0.32 + i * 0.46
    bg = ACCENT1 if i % 2 == 0 else RGBColor(0x10,0x18,0x28)
    add_rect(s3, 0.25, row_y, 7.0, 0.44, bg)
    tcol = tier_cols[tier]
    mcol = model_cols[model]
    add_text(s3, agent, 0.3,  row_y+0.05, 2.5,  0.35, font_size=11, color=TEAL)
    add_text(s3, tier,  2.85, row_y+0.05, 0.45, 0.35, font_size=11, bold=True, color=tcol, align=PP_ALIGN.CENTER)
    add_text(s3, model, 3.4,  row_y+0.05, 1.15, 0.35, font_size=11, color=mcol)
    add_text(s3, role,  4.65, row_y+0.05, 2.55, 0.35, font_size=11, color=LIGHT_GRAY)

# Right: orchestration patterns
add_text(s3, "Orchestration Patterns", 7.6, 0.85, 5.5, 0.45,
         font_size=17, bold=True, color=ORANGE)

patterns = [
    ("Pipeline",      TEAL,
     "/geos:fluids → /geos → /geos:validate\n"
     "Output flows sequentially through agents"),
    ("Fan-out",       GREEN,
     "/geos dispatches mesh + fluids in parallel\n"
     "Results merged before XML assembly"),
    ("Feedback Loop", PURPLE,
     "Produce → Review → Re-invoke (bounded)\n"
     "e.g. XML → validate → fix → re-validate"),
]

py = 1.38
for name, col, desc in patterns:
    add_rect(s3, 7.55, py, 5.55, 1.5, ACCENT1)
    add_rect(s3, 7.55, py, 0.07, 1.5, col)
    add_text(s3, name, 7.72, py+0.08, 5.2, 0.4,
             font_size=14, bold=True, color=col)
    tb = s3.shapes.add_textbox(Inches(7.72), Inches(py+0.5),
                                Inches(5.2), Inches(0.9))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(desc.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run()
        r.text = line
        r.font.size = Pt(13)
        r.font.color.rgb = LIGHT_GRAY
    py += 1.62

# Cognitive tier legend
add_rect(s3, 7.55, py+0.1, 5.55, 1.5, ACCENT1)
add_text(s3, "Cognitive Tiers  →  Model Routing",
         7.7, py+0.18, 5.3, 0.35, font_size=13, bold=True, color=ORANGE)
tiers = [
    ("Tier 1  Haiku",   "Lookup, extraction, formatting",          RGBColor(0x42,0xA5,0xF5)),
    ("Tier 2  Sonnet",  "Synthesis, assembly, validation",         ORANGE),
    ("Tier 3  Opus",    "Planning, multi-step reasoning, dialogue", GREEN),
]
ty = py + 0.55
for label, desc, col in tiers:
    add_rect(s3, 7.62, ty, 1.5, 0.28, col)
    add_text(s3, label, 7.63, ty+0.02, 1.48, 0.24,
             font_size=10, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)
    add_text(s3, desc, 9.2, ty+0.02, 3.8, 0.24, font_size=10, color=LIGHT_GRAY)
    ty += 0.33

add_text(s3, "Key rule: all agent-to-agent handoffs pass structured data, never free-form prose",
         7.62, ty+0.12, 5.4, 0.4, font_size=11, italic=True, color=RGBColor(0x90,0x90,0x90))

add_text(s3, "3 / 4", 12.8, 7.1, 0.4, 0.3, font_size=11,
         color=LIGHT_GRAY, align=PP_ALIGN.RIGHT)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — End-to-End Workflow & Knowledge Base
# ════════════════════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(blank)
set_bg(s4, DARK_BG)
add_rect(s4, 0, 0, 13.33, 0.12, TEAL)
add_rect(s4, 0, 7.38, 13.33, 0.12, TEAL)
add_rect(s4, 0, 0.12, 0.08, 7.26, ORANGE)

add_text(s4, "End-to-End Workflow  &  Knowledge Base",
         0.25, 0.18, 13, 0.55, font_size=28, bold=True, color=TEAL)

# ── Workflow arrow chain ────────────────────────────────────────────────────
steps = [
    ("Natural\nLanguage\nInput",  TEAL,   0.3),
    ("Mesh\n/geos:mesh",          GREEN,  2.4),
    ("Fluids\n/geos:fluids",      ORANGE, 4.5),
    ("XML Assembly\n/geos",       PURPLE, 6.6),
    ("Validation\n/geos:validate",RED_SOFT,8.7),
    ("Post-Process\n/geos:post",  RGBColor(0xFF,0xB7,0x4D), 10.8),
]

box_w = 1.95
box_h = 1.4
box_y = 0.9

for i, (label, col, lft) in enumerate(steps):
    add_rect(s4, lft, box_y, box_w, box_h, col)
    add_text(s4, label, lft+0.05, box_y+0.1, box_w-0.1, box_h-0.15,
             font_size=13, bold=(i == 0), color=DARK_BG, align=PP_ALIGN.CENTER)
    if i < len(steps) - 1:
        arr_x = lft + box_w + 0.02
        add_text(s4, "▶", arr_x, box_y + 0.45, 0.32, 0.5,
                 font_size=20, bold=True, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# ── Output box ─────────────────────────────────────────────────────────────
add_rect(s4, 0.3, 2.55, 12.7, 0.05, ORANGE)
add_text(s4, "Output: validated GEOS XML  (200-500 lines with correct cross-references)",
         0.3, 2.65, 12.7, 0.45, font_size=15, italic=True, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# ── Knowledge modules ──────────────────────────────────────────────────────
add_text(s4, "Knowledge Base  —  Domain Expertise Encoded in 5 Modules",
         0.3, 3.2, 12.7, 0.45, font_size=17, bold=True, color=ORANGE)

km = [
    ("field_names.py",    TEAL,
     "Solver type → valid BC / IC field names\nPrevents wrong field assignments"),
    ("cross_refs.py",     GREEN,
     "Attribute → section mapping graph\nEnsures all name references resolve"),
    ("sanity_rules.py",   ORANGE,
     "Physics heuristics + structural checks\nCatches unphysical parameter combos"),
    ("fluid_models.py",   PURPLE,
     "NL keywords → solver + constitutive\nassembly for 6 fluid scenarios"),
    ("lessons_learned.md",RED_SOFT,
     "Runtime error patterns and fixes\nGrows via /geos:curate-errors"),
]

km_w = 2.45
km_h = 2.65
km_y = 3.75
km_start = 0.28

for i, (name, col, desc) in enumerate(km):
    lft = km_start + i * (km_w + 0.07)
    add_rect(s4, lft, km_y, km_w, km_h, ACCENT1)
    add_rect(s4, lft, km_y, km_w, 0.07, col)
    add_text(s4, name, lft+0.1, km_y+0.12, km_w-0.15, 0.45,
             font_size=12, bold=True, color=col)
    tb = s4.shapes.add_textbox(Inches(lft+0.1), Inches(km_y+0.6),
                                Inches(km_w-0.2), Inches(km_h-0.7))
    tf = tb.text_frame
    tf.word_wrap = True
    for j, line in enumerate(desc.split("\n")):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.space_before = Pt(3)
        r = p.add_run()
        r.text = line
        r.font.size = Pt(12)
        r.font.color.rgb = LIGHT_GRAY

# Key principle callout
add_rect(s4, 0.28, 6.5, 12.72, 0.72, RGBColor(0x0D, 0x47, 0x55))
add_text(s4,
         "Key principle:  Tools never hardcode domain logic — they read from knowledge modules. "
         "New patterns discovered at runtime are added to the appropriate module, not to tools.",
         0.45, 6.55, 12.4, 0.6, font_size=14, italic=True, color=TEAL)

add_text(s4, "4 / 4", 12.8, 7.1, 0.4, 0.3, font_size=11,
         color=LIGHT_GRAY, align=PP_ALIGN.RIGHT)


# ── Save ─────────────────────────────────────────────────────────────────────
out = "/home/user/agents4geos/agents4geos_slides.pptx"
prs.save(out)
print(f"Saved: {out}")
